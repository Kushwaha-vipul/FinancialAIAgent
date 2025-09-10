import os
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter


REPORTS_DIR = "reports"
os.makedirs(REPORTS_DIR, exist_ok=True)


def generate_ppt(title: str, content: str, filename: str = "report.pptx"):
    file_path = os.path.join(REPORTS_DIR, filename)

    prs = Presentation()
    

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Details"
    slide.placeholders[1].text = content

    prs.save(file_path)
    return file_path  


def generate_pdf(title: str, content: str, filename: str = "report.pdf"):
    file_path = os.path.join(REPORTS_DIR, filename)

    c = canvas.Canvas(file_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, height - 50, title)
    
    c.setFont("Helvetica", 12)
    text = c.beginText(50, height - 100)
    for line in content.split("\n"):
        text.textLine(line)
    c.drawText(text)
    
    c.showPage()
    c.save()
    return file_path
